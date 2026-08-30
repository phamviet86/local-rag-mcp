from __future__ import annotations

import mimetypes
from pathlib import Path
from typing import Any

from .models import ExtractedDocument, SourceSpan
from .ocr_runtime import OCRRuntimeManager


class Extractor:
    def __init__(
        self, ocr_runtime: OCRRuntimeManager, *, ocr_enabled: bool = True, ocr_offline: bool = False
    ):
        self.ocr_runtime = ocr_runtime
        self.ocr_enabled = ocr_enabled
        self.ocr_offline = ocr_offline

    def extract(self, path: Path) -> ExtractedDocument:
        suffix = path.suffix.lower()
        if suffix in {".txt", ".md", ".markdown"}:
            return _text(path, markdown=suffix != ".txt")
        if suffix == ".docx":
            return _docx(path)
        if suffix == ".xlsx":
            return _xlsx(path)
        if suffix == ".pptx":
            return _pptx(path)
        if suffix == ".pdf":
            return self._pdf(path)
        raise ValueError(f"unsupported file type: {suffix}")

    def _pdf(self, path: Path) -> ExtractedDocument:
        import pdf_inspector

        classification = pdf_inspector.process_pdf(str(path))
        native = pdf_inspector.extract_pages_markdown(str(path))
        routed = list(classification.pages_needing_ocr)
        pages: dict[int, dict[str, Any]] = {}
        for native_page in native.pages:
            number = int(native_page.page) + 1
            markdown = (native_page.markdown or "").strip()
            if markdown:
                pages[number] = {"text": markdown, "source": "native"}
        reviews: list[dict[str, Any]] = []
        if routed:
            if self.ocr_enabled and self.ocr_runtime.configure():
                try:
                    model_directory = self.ocr_runtime.resolved_model_dir()
                    result = pdf_inspector.process_pdf_with_ocr(
                        str(path),
                        page_numbers=routed,
                        model_directory=str(model_directory) if model_directory else None,
                        offline=self.ocr_offline,
                    )
                    for ocr_page in result.pages:
                        number = int(ocr_page.page_number)
                        markdown = (ocr_page.markdown or "").strip()
                        provenance = ocr_page.provenance
                        if markdown:
                            pages[number] = {
                                "text": markdown,
                                "source": provenance.source,
                                "confidence": provenance.ocr_confidence,
                            }
                        if not markdown or provenance.hosted_recommended:
                            reviews.append(
                                {
                                    "page": number,
                                    "reason": "low_quality_ocr"
                                    if markdown
                                    else "blank_or_failed_ocr",
                                    "detail": {"warnings": list(provenance.warnings)},
                                }
                            )
                except Exception as exc:
                    reviews.extend(
                        {"page": page, "reason": "ocr_failed", "detail": {"error": str(exc)}}
                        for page in routed
                    )
            else:
                reviews.extend(
                    {
                        "page": page,
                        "reason": "ocr_runtime_missing",
                        "detail": {
                            "action": "run local-rag-mcp setup --full, then reindex --reextract"
                        },
                    }
                    for page in routed
                )
        complex_pages = set(classification.pages_with_tables) | set(
            classification.pages_with_columns
        )
        reviews.extend(
            {"page": page, "reason": "complex_layout", "detail": {}} for page in complex_pages
        )
        builder = _Builder()
        for number, page_data in sorted(pages.items()):
            builder.add(
                page_data["text"],
                "pdf_page",
                f"page:{number}",
                {key: value for key, value in page_data.items() if key != "text"},
            )
        metadata = {
            "pdf_type": classification.pdf_type,
            "page_count": classification.page_count,
            "classification_confidence": classification.confidence,
            "pages_routed_to_ocr": routed,
            "complex_layout": bool(classification.is_complex_layout),
        }
        return ExtractedDocument(builder.text, builder.spans, "application/pdf", metadata, reviews)


class _Builder:
    def __init__(self) -> None:
        self.parts: list[str] = []
        self.spans: list[SourceSpan] = []
        self.length = 0

    def add(
        self, value: str, kind: str, locator: str, metadata: dict[str, Any] | None = None
    ) -> None:
        cleaned = value.strip()
        if not cleaned:
            return
        separator = "\n\n" if self.parts else ""
        self.parts.append(separator + cleaned)
        start = self.length + len(separator)
        self.length += len(separator) + len(cleaned)
        self.spans.append(SourceSpan(kind, locator, start, self.length, metadata or {}))

    @property
    def text(self) -> str:
        return "".join(self.parts)


def _text(path: Path, markdown: bool) -> ExtractedDocument:
    text = path.read_text(encoding="utf-8", errors="replace")
    builder = _Builder()
    for number, line in enumerate(text.splitlines(), 1):
        if line.strip():
            kind = "markdown_line" if markdown else "text_line"
            builder.add(line, kind, f"line:{number}")
    media = "text/markdown" if markdown else "text/plain"
    return ExtractedDocument(builder.text, builder.spans, media)


def _docx(path: Path) -> ExtractedDocument:
    from docx import Document

    document, builder = Document(str(path)), _Builder()
    for index, paragraph in enumerate(document.paragraphs, 1):
        builder.add(paragraph.text, "docx_paragraph", f"paragraph:{index}")
    for table_index, table in enumerate(document.tables, 1):
        for row_index, row in enumerate(table.rows, 1):
            for column_index, cell in enumerate(row.cells, 1):
                builder.add(
                    cell.text,
                    "docx_cell",
                    f"table:{table_index}/row:{row_index}/column:{column_index}",
                )
    return ExtractedDocument(builder.text, builder.spans, _media(path))


def _xlsx(path: Path) -> ExtractedDocument:
    from openpyxl import load_workbook

    workbook, builder = load_workbook(path, read_only=True, data_only=True), _Builder()
    try:
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value is not None:
                        builder.add(
                            str(cell.value),
                            "xlsx_cell",
                            f"sheet:{sheet.title}/cell:{cell.coordinate}",
                        )
    finally:
        workbook.close()
    return ExtractedDocument(
        builder.text, builder.spans, _media(path), {"sheets": workbook.sheetnames}
    )


def _pptx(path: Path) -> ExtractedDocument:
    from pptx import Presentation

    presentation, builder = Presentation(str(path)), _Builder()
    for slide_index, slide in enumerate(presentation.slides, 1):
        for shape_index, shape in enumerate(slide.shapes, 1):
            text = getattr(shape, "text", "")
            builder.add(text, "pptx_shape", f"slide:{slide_index}/shape:{shape_index}")
    return ExtractedDocument(
        builder.text, builder.spans, _media(path), {"slides": len(presentation.slides)}
    )


def _media(path: Path) -> str:
    return mimetypes.guess_type(path.name)[0] or "application/octet-stream"
