import tempfile
import unittest
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from pptx import Presentation

from local_rag.extract import Extractor
from local_rag.ocr_runtime import OCRRuntimeManager
from tests.helpers import write_text_pdf


class ExtractionTests(unittest.TestCase):
    def setUp(self):
        self.temp = tempfile.TemporaryDirectory()
        self.base = Path(self.temp.name)
        self.extractor = Extractor(OCRRuntimeManager(self.base / "runtime", self.base / "models"))

    def tearDown(self):
        self.temp.cleanup()

    def test_text_markdown_and_office_provenance(self):
        text = self.base / "note.txt"
        text.write_text("first\nsecond", encoding="utf-8")
        self.assertEqual(self.extractor.extract(text).spans[1].locator, "line:2")

        markdown = self.base / "note.md"
        markdown.write_text("# Heading\nbody", encoding="utf-8")
        self.assertEqual(self.extractor.extract(markdown).spans[0].kind, "markdown_line")

        docx = self.base / "sample.docx"
        document = Document()
        document.add_paragraph("DOCX paragraph")
        document.add_table(rows=1, cols=1).cell(0, 0).text = "DOCX cell"
        document.save(docx)
        extracted = self.extractor.extract(docx)
        self.assertIn("DOCX paragraph", extracted.text)
        self.assertTrue(any(span.kind == "docx_cell" for span in extracted.spans))

        xlsx = self.base / "sample.xlsx"
        workbook = Workbook()
        workbook.active.title = "Data"
        workbook.active["B2"] = "XLSX value"
        workbook.save(xlsx)
        extracted = self.extractor.extract(xlsx)
        self.assertEqual(extracted.spans[0].locator, "sheet:Data/cell:B2")

        pptx = self.base / "sample.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "PPTX title"
        presentation.save(pptx)
        extracted = self.extractor.extract(pptx)
        self.assertIn("PPTX title", extracted.text)
        self.assertTrue(extracted.spans[0].locator.startswith("slide:1/shape:"))

    def test_pdf_inspector_native_extraction_and_classification(self):
        pdf = self.base / "native.pdf"
        write_text_pdf(pdf, "Positioned native content")
        extracted = self.extractor.extract(pdf)
        self.assertIn("Positioned native content", extracted.text)
        self.assertEqual(extracted.metadata["pdf_type"], "text_based")
        self.assertEqual(extracted.spans[0].locator, "page:1")
        self.assertEqual(extracted.metadata["pages_routed_to_ocr"], [1])
        self.assertEqual(extracted.reviews[0]["reason"], "ocr_runtime_missing")


if __name__ == "__main__":
    unittest.main()
